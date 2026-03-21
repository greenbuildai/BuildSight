from pptx import Presentation

prs = Presentation(r"E:\Company\Green Build AI\Prototypes\BuildSight\Batch 12 1st_review(Final).pptx")

with open("pptx_text.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides):
        if i >= 3:
            break
        f.write(f"\n--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                f.write(shape.text + "\n")
